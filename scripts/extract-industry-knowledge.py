from pathlib import Path
import json
import re
import sys
from datetime import datetime, timezone
from pypdf import PdfReader
from pptx import Presentation

sys.stdout.reconfigure(encoding="utf-8")

ROOT = Path("knowledge/industry")
OUTPUT_DIR = ROOT / "processed"
MANIFEST_PATH = OUTPUT_DIR / "industry-manifest.json"

PROFILES = {
    "产品认证与合规 Esther.pdf": ("certification-compliance", 4, "产品认证与合规"),
    "产品通用功能_Ivy.pptx": ("network-product-capabilities", 3, "网络产品通用功能"),
    "认知启航：网络是怎么一回事儿-Reed20260707.pdf": ("network-access-topology", 4, "网络接入、覆盖与拓扑"),
    "市场操盘及渠道深耕-Jacky.pptx": ("channel-go-to-market", 3, "市场操盘与渠道深耕"),
    "网络产品品牌与行业研究报告.pptx": ("brand-market-landscape", 2, "网络品牌与区域市场研究"),
    "网络基础知识_Sam.pdf": ("network-foundations", 4, "网络基础知识"),
}


def clean_text(value):
    text = str(value or "").replace("\x00", "")
    text = re.sub(r"[\t ]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def pdf_units(path):
    reader = PdfReader(path)
    return [(index, clean_text(page.extract_text() or "")) for index, page in enumerate(reader.pages, start=1)]


def pptx_units(path):
    deck = Presentation(path)
    units = []
    for index, slide in enumerate(deck.slides, start=1):
        parts = []
        links = set()
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                for paragraph in shape.text_frame.paragraphs:
                    text = clean_text(paragraph.text)
                    if text:
                        parts.append(text)
                    for run in paragraph.runs:
                        address = run.hyperlink.address
                        if address:
                            links.add(address)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    parts.append(" | ".join(clean_text(cell.text) for cell in row.cells))
        if links:
            parts.append("Sources: " + ", ".join(sorted(links)))
        units.append((index, clean_text("\n".join(parts))))
    return units


def main():
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    documents = []
    for filename, (topic, authority, title) in PROFILES.items():
        path = ROOT / filename
        if not path.exists():
            continue
        units = pdf_units(path) if path.suffix.lower() == ".pdf" else pptx_units(path)
        unit_label = "page" if path.suffix.lower() == ".pdf" else "slide"
        body = []
        for index, text in units:
            if text:
                body.append(f"## Source {unit_label} {index}\n\n{text}")
            else:
                body.append(f"## Source {unit_label} {index}\n\n[Visual-only or no extractable text]")
        markdown = "\n\n".join([
            f"# {title}",
            "## Source metadata",
            f"- Source file: {filename}\n- Topic: {topic}\n- Material type: internal training/research\n- Authority level: {authority}/5\n- Unit count: {len(units)}",
            *body,
        ]) + "\n"
        output_path = OUTPUT_DIR / f"{topic}.md"
        output_path.write_text(markdown, encoding="utf-8")
        modified = datetime.fromtimestamp(path.stat().st_mtime, tz=timezone.utc).isoformat()
        documents.append({
            "externalId": f"industry:{topic}",
            "title": title,
            "topic": topic,
            "authorityLevel": authority,
            "language": "zh-CN",
            "sourceFile": filename,
            "sourceType": "internal-training-material",
            "knowledgeFile": output_path.as_posix(),
            "unitCount": len(units),
            "capturedAt": modified,
        })
    MANIFEST_PATH.write_text(json.dumps({"documents": documents}, ensure_ascii=False, indent=2), encoding="utf-8")
    print(json.dumps({"documents": len(documents), "units": sum(item["unitCount"] for item in documents), "manifest": str(MANIFEST_PATH)}, ensure_ascii=False))


if __name__ == "__main__":
    main()
