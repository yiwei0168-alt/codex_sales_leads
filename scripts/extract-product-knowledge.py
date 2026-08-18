from pathlib import Path
import hashlib
import json
import re
import sys
from openpyxl import load_workbook
from pypdf import PdfReader
from pptx import Presentation

ROOT = Path("knowledge/product")
PROCESSED_DIR = ROOT / "processed"
DATASHEET_OUTPUT_DIR = PROCESSED_DIR / "datasheets"
CATALOG_OUTPUT_DIR = PROCESSED_DIR / "catalog"
REFERENCE_OUTPUT_DIR = PROCESSED_DIR / "references"
CATALOG_PATH = PROCESSED_DIR / "product-catalog.json"


def clean_text(value):
    text = str(value or "").replace("\x00", "")
    text = text.replace("¡Á", "×").replace("Â×", "×")
    text = re.sub(r"[\t ]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.encode("utf-8", "replace").decode("utf-8").strip()


def safe_slug(value):
    slug = re.sub(r"[^a-z0-9]+", "-", value.lower()).strip("-")
    return slug[:80] or hashlib.sha256(value.encode("utf-8")).hexdigest()[:16]


def read_catalog():
    workbook = load_workbook(ROOT / "Cudy products list.xlsx", read_only=True, data_only=True)
    sheet = workbook["Price List"]
    category = None
    products = []
    for row in sheet.iter_rows(values_only=True):
        values = [clean_text(value) for value in row]
        first, second, third = (values + ["", "", ""])[:3]
        if not first or first in {"Quotation", "Model No."}:
            continue
        if first and not second and not third:
            category = first
            continue
        if not category or not second:
            continue
        products.append({
            "model": first,
            "productName": second,
            "category": category,
            "description": third,
            "brand": "Cudy Technology",
            "lifecycleStatus": "unknown",
            "sourceFile": "Cudy products list.xlsx",
        })
    return products


def matched_models(path, products):
    stem = path.stem.upper()
    matches = []
    for product in sorted(products, key=lambda item: len(item["model"]), reverse=True):
        model = product["model"].upper()
        if re.search(rf"(?<![A-Z0-9]){re.escape(model)}(?![A-Z0-9])", stem):
            matches.append(product["model"])
    return matches


def datasheet_version(path):
    match = re.search(r"(?:^|[_-])V?(\d+(?:\.\d+)+)(?:[_-]|$)", path.stem, re.I)
    return f"V{match.group(1)}" if match else None


def extract_pdf(path):
    reader = PdfReader(path)
    pages = []
    for index, page in enumerate(reader.pages, start=1):
        text = clean_text(page.extract_text() or "")
        if text:
            pages.append(f"## Page {index}\n\n{text}")
    return "\n\n".join(pages), len(reader.pages)


def extract_pptx(path):
    presentation = Presentation(path)
    slides = []
    for index, slide in enumerate(presentation.slides, start=1):
        parts = []
        for shape in slide.shapes:
            text = clean_text(getattr(shape, "text", ""))
            if text:
                parts.append(text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    line = " | ".join(clean_text(cell.text) for cell in row.cells)
                    if line.strip(" |"): parts.append(line)
        if parts:
            slides.append(f"## Slide {index}\n\n" + "\n\n".join(parts))
    return "\n\n".join(slides), len(presentation.slides)


def write_document(output_path, title, identity, content):
    output_path.parent.mkdir(parents=True, exist_ok=True)
    output_path.write_text("\n\n".join([f"# {title}", identity, content or "No extractable text."]) + "\n", encoding="utf-8")


def extract_datasheets(products):
    records = []
    paths = [path for path in ROOT.rglob("*.pdf") if PROCESSED_DIR not in path.parents and path.parent != ROOT and "datasheet" in path.name.lower()]
    for path in sorted(paths):
        models = matched_models(path, products)
        version = datasheet_version(path)
        content, page_count = extract_pdf(path)
        category = path.parent.name
        identity = "\n".join([
            "## Product identity",
            "- Brand: Cudy Technology",
            f"- Models: {', '.join(models) if models else 'Unmatched'}",
            f"- Category: {category}",
            f"- Datasheet version: {version or 'Unknown'}",
            f"- Source file: {path.relative_to(ROOT).as_posix()}",
            f"- Source pages: {page_count}",
        ])
        output_path = DATASHEET_OUTPUT_DIR / safe_slug(category) / f"{safe_slug(path.stem)}.md"
        title_models = " / ".join(models) if models else path.stem
        write_document(output_path, f"{title_models} Datasheet", identity, content)
        records.append({
            "model": models[0] if models else path.stem.split("_")[0],
            "relatedModels": models,
            "productName": f"{title_models} Datasheet",
            "category": category,
            "description": "",
            "brand": "Cudy Technology",
            "lifecycleStatus": "unknown",
            "sourceFile": path.relative_to(ROOT).as_posix(),
            "datasheetVersion": version,
            "datasheetFile": path.name,
            "knowledgeFile": output_path.as_posix(),
            "pageCount": page_count,
        })
    return records


def extract_catalog_documents(products):
    records = []
    categories = sorted({product["category"] for product in products})
    for category in categories:
        items = [product for product in products if product["category"] == category]
        lines = [f"## {item['model']} - {item['productName']}\n\n{item['description'] or 'No catalog description.'}" for item in items]
        output_path = CATALOG_OUTPUT_DIR / f"{safe_slug(category)}.md"
        identity = "\n".join([
            "## Catalog identity", "- Brand: Cudy Technology", f"- Category: {category}",
            f"- Models: {len(items)}", "- Source file: Cudy products list.xlsx",
        ])
        write_document(output_path, f"Cudy {category} Product Catalog", identity, "\n\n".join(lines))
        records.append({
            "title": f"Cudy {category} Product Catalog", "sourceFile": "Cudy products list.xlsx",
            "knowledgeFile": output_path.as_posix(), "sourceType": "product-catalog-category",
            "category": category, "relatedModels": [item["model"] for item in items],
        })
    return records


def extract_reference_documents(products):
    records = []
    paths = sorted([*ROOT.glob("*.pdf"), *ROOT.glob("*.pptx")])
    for path in paths:
        if path.suffix.lower() == ".pdf":
            content, unit_count = extract_pdf(path)
            unit_name = "pageCount"
        else:
            content, unit_count = extract_pptx(path)
            unit_name = "slideCount"
        models = matched_models(path, products)
        output_path = REFERENCE_OUTPUT_DIR / f"{safe_slug(path.stem)}.md"
        identity = "\n".join([
            "## Reference identity", "- Brand: Cudy Technology",
            f"- Source file: {path.name}", f"- Related models in filename: {', '.join(models) if models else 'Not specified'}",
            f"- {unit_name}: {unit_count}",
        ])
        write_document(output_path, path.stem, identity, content)
        records.append({
            "title": path.stem, "sourceFile": path.name, "knowledgeFile": output_path.as_posix(),
            "sourceType": "product-training-reference", "relatedModels": models, unit_name: unit_count,
        })
    return records


def main():
    products = read_catalog()
    datasheets = extract_datasheets(products)
    catalog_documents = extract_catalog_documents(products)
    references = extract_reference_documents(products)
    CATALOG_PATH.parent.mkdir(parents=True, exist_ok=True)
    CATALOG_PATH.write_text(json.dumps({
        "allProducts": products,
        "datasheets": datasheets,
        "wifiRouters": [item for item in datasheets if item["category"] == "Wi-Fi Router"],
        "catalogDocuments": catalog_documents,
        "references": references,
    }, ensure_ascii=False, indent=2), encoding="utf-8")
    print(json.dumps({
        "catalogProducts": len(products), "categories": len(catalog_documents),
        "datasheets": len(datasheets), "references": len(references), "output": str(CATALOG_PATH),
    }, ensure_ascii=False))


if __name__ == "__main__":
    try:
        main()
    except Exception as error:
        print(f"Product extraction failed: {error}", file=sys.stderr)
        raise
